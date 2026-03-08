"""Tests for create_pptx tool."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from office_agent.tools.pptx_tool import create_pptx


def test_create_pptx_basic(sample_pptx_plan: dict, tmp_path: Path) -> None:
    result = create_pptx(plan=sample_pptx_plan, out_dir=str(tmp_path))
    assert result["tool_used"] == "create_pptx"
    out = Path(result["output_path"])
    assert out.exists()
    assert out.suffix == ".pptx"


def test_create_pptx_slide_count(sample_pptx_plan: dict, tmp_path: Path) -> None:
    result = create_pptx(plan=sample_pptx_plan, out_dir=str(tmp_path))
    prs = Presentation(result["output_path"])
    assert len(prs.slides) == len(sample_pptx_plan["slides"])


def test_create_pptx_title_slide(sample_pptx_plan: dict, tmp_path: Path) -> None:
    result = create_pptx(plan=sample_pptx_plan, out_dir=str(tmp_path))
    prs = Presentation(result["output_path"])
    first_slide = prs.slides[0]
    title_shape = first_slide.shapes.title
    assert title_shape is not None
    assert sample_pptx_plan["slides"][0]["title"] in title_shape.text


def test_create_pptx_notes(sample_pptx_plan: dict, tmp_path: Path) -> None:
    result = create_pptx(plan=sample_pptx_plan, out_dir=str(tmp_path))
    prs = Presentation(result["output_path"])
    # Slide index 1 has notes
    notes_text = prs.slides[1].notes_slide.notes_text_frame.text
    expected = sample_pptx_plan["slides"][1]["notes"]
    if expected:
        assert expected in notes_text


def test_tool_registry_pptx(registry, sample_pptx_plan: dict, tmp_path: Path) -> None:
    result = registry.call("create_pptx", plan=sample_pptx_plan, out_dir=str(tmp_path))
    assert result["tool_used"] == "create_pptx"
