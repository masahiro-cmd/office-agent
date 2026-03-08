"""create_pptx: Generate a PowerPoint presentation from a plan JSON."""

from __future__ import annotations

import logging
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

logger = logging.getLogger(__name__)


def create_pptx(plan: dict, out_dir: str, template_dir: str = "./templates") -> dict:
    """
    Generate a .pptx file from a validated plan dict.

    Args:
        plan: Validated presentation plan (must conform to pptx_schema.json).
        out_dir: Directory where the output file will be saved.
        template_dir: Directory containing optional .pptx template files.

    Returns:
        {"output_path": str, "tool_used": "create_pptx"}
    """
    title: str = plan.get("presentation_title", "Untitled")
    slides_data: list[dict] = plan.get("slides", [])
    template_setting: str = plan.get("template", "default")

    # Load template if specified
    template_path = _resolve_template(template_setting, template_dir)
    if template_path:
        logger.info(f"Using template: {template_path}")
        prs = Presentation(str(template_path))
    else:
        prs = Presentation()

    for slide_plan in slides_data:
        _add_slide(prs, slide_plan)

    out_path = Path(out_dir)
    out_path.mkdir(parents=True, exist_ok=True)
    safe_title = _safe_filename(title)
    file_path = out_path / f"{safe_title}.pptx"
    prs.save(str(file_path))
    logger.info(f"Saved pptx: {file_path}")

    return {"output_path": str(file_path), "tool_used": "create_pptx"}


def _add_slide(prs: Presentation, slide_plan: dict) -> None:
    slide_type: str = slide_plan.get("slide_type", "content")
    slide_title: str = slide_plan.get("title", "")
    bullets: list[str] = slide_plan.get("bullets") or []
    table_data: dict | None = slide_plan.get("table")
    notes_text: str = slide_plan.get("notes", "")

    # Choose layout
    layout_idx = 0 if slide_type == "title" else 1
    try:
        layout = prs.slide_layouts[layout_idx]
    except IndexError:
        layout = prs.slide_layouts[0]

    slide = prs.slides.add_slide(layout)

    # Set title placeholder
    if slide.shapes.title:
        slide.shapes.title.text = slide_title

    # Set body/content placeholder with bullets
    if bullets and len(slide.placeholders) > 1:
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()
        for i, bullet in enumerate(bullets):
            if i == 0:
                para = tf.paragraphs[0]
            else:
                para = tf.add_paragraph()
            para.text = bullet
            para.level = 0

    # Add table if provided
    if table_data:
        headers = table_data.get("headers", [])
        rows = table_data.get("rows", [])
        if headers:
            rows_count = 1 + len(rows)
            cols_count = len(headers)
            left = Inches(1)
            top = Inches(3)
            width = Inches(8)
            height = Inches(0.5 * rows_count)
            table_shape = slide.shapes.add_table(rows_count, cols_count, left, top, width, height)
            tbl = table_shape.table
            for c_idx, h in enumerate(headers):
                tbl.cell(0, c_idx).text = str(h)
            for r_idx, row in enumerate(rows):
                for c_idx, val in enumerate(row):
                    tbl.cell(r_idx + 1, c_idx).text = str(val)

    # Speaker notes
    if notes_text:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes_text


def _resolve_template(template_setting: str, template_dir: str) -> Path | None:
    if template_setting == "default":
        td = Path(template_dir) / "pptx"
        if td.exists():
            matches = list(td.glob("*.pptx"))
            if matches:
                return matches[0]
        return None
    # Treat as explicit path
    p = Path(template_setting)
    if p.exists():
        return p
    return None


def _safe_filename(name: str) -> str:
    invalid = r'\/:*?"<>|'
    for ch in invalid:
        name = name.replace(ch, "_")
    return name[:100]
